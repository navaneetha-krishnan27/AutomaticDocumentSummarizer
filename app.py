import os
import io
import re
import json
import time
from collections import Counter, defaultdict
from datetime import datetime
from functools import wraps

from flask import Flask, render_template, request, redirect, url_for, flash, Response, jsonify, send_file, make_response
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, login_required, logout_user, current_user
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash
import PyPDF2
from docx import Document
import matplotlib
matplotlib.use('Agg')
from fpdf import FPDF
import spacy
from deep_translator import GoogleTranslator
import zipfile
from io import BytesIO
from PIL import Image

# --- ML Imports ---
from sentence_transformers import SentenceTransformer, util
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# --- Sumy Imports ---
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer 
from sumy.summarizers.lex_rank import LexRankSummarizer
from sumy.utils import get_stop_words

# --- PPTX Import ---
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

# Load Spacy Model
try:
    nlp = spacy.load("en_core_web_sm")
except:
    import sys
    print("Spacy model not found. Run: python -m spacy download en_core_web_sm")
    sys.exit(1)

# --- Load ML Model (Offline Compatible) ---
model_path = './local_ml_model'
if os.path.exists(model_path):
    print("Loading ML model from local storage...")
    ml_model = SentenceTransformer(model_path)
else:
    print("Local model not found. Attempting to download...")
    try:
        ml_model = SentenceTransformer('all-MiniLM-L6-v2')
    except Exception as e:
        print(f"FAILED to load ML Model: {e}")
        ml_model = None

app = Flask(__name__)
app.config['SECRET_KEY'] = 'secure-key-123'

# --- DATABASE CONFIG ---
app.config['SQLALCHEMY_DATABASE_URI'] = 'mysql+pymysql://root:password@local/nlp_fresh_db'
app.config['UPLOAD_FOLDER'] = 'uploads'
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Initialize DB and Login Manager
db = SQLAlchemy(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

# --- MODELS ---

class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(150), unique=True, nullable=False)
    email = db.Column(db.String(150), unique=True, nullable=False)
    password = db.Column(db.String(150), nullable=False)
    dob = db.Column(db.String(20), nullable=False)
    role = db.Column(db.String(20), default='user') 
    is_blocked = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class Summary(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(150))
    original_text = db.Column(db.Text)
    summary_text = db.Column(db.Text)
    language = db.Column(db.String(10), default='en')
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)

class ChatMessage(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    summary_id = db.Column(db.Integer, db.ForeignKey('summary.id'), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    message = db.Column(db.Text, nullable=False)
    sender = db.Column(db.String(10), nullable=False)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

class ActivityLog(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True)
    action = db.Column(db.String(50), nullable=False)
    details = db.Column(db.String(255))
    ip_address = db.Column(db.String(50))
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

class SiteSetting(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    key_name = db.Column(db.String(50), unique=True)
    value = db.Column(db.String(255))

class Rating(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'))
    stars = db.Column(db.Integer)
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)

class Notification(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    message = db.Column(db.Text, nullable=False)
    type = db.Column(db.String(20), default='info') 
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

# --- HELPERS ---

def log_activity(action, details=None):
    uid = current_user.id if current_user.is_authenticated else None
    ip = request.remote_addr
    new_log = ActivityLog(user_id=uid, action=action, details=details, ip_address=ip)
    db.session.add(new_log)
    db.session.commit()

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated or current_user.role != 'admin':
            flash("🚫 Access denied. Admins only.")
            return redirect(url_for('dashboard'))
        return f(*args, **kwargs)
    return decorated_function

def is_maintenance_active():
    maint = SiteSetting.query.filter_by(key_name='maintenance_mode').first()
    if maint and maint.value == 'true':
        if current_user.is_authenticated and current_user.role == 'admin':
            return False
        return True
    return False

# --- CONTEXT PROCESSOR ---
@app.context_processor
def inject_global_vars():
    # Pass Rating Setting
    r_setting = SiteSetting.query.filter_by(key_name='enable_ratings').first()
    rating_enabled = True if r_setting and r_setting.value == 'true' else False
    
    # Check if the specific user has ALREADY rated in the CURRENT session
    user_has_rated = False
    if current_user.is_authenticated:
        # Get the time when the admin last enabled ratings
        start_setting = SiteSetting.query.filter_by(key_name='rating_start_time').first()
        start_time = datetime.min # Default to very old if not set
        
        if start_setting and start_setting.value:
            try:
                start_time = datetime.fromisoformat(start_setting.value)
            except ValueError:
                pass
        
        # Check if user has a rating AFTER that start time
        if Rating.query.filter(Rating.user_id == current_user.id, Rating.timestamp >= start_time).first():
            user_has_rated = True

    # Pass Latest Announcement
    announcement = Notification.query.order_by(Notification.created_at.desc()).first()
    active_announcement = announcement.message if announcement else None

    return dict(
        rating_enabled=rating_enabled, 
        active_announcement=active_announcement,
        user_has_rated=user_has_rated
    )

# --- UTILITY FUNCTIONS ---

def clean_text(text):
    # 1. Fix Broken Words
    text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)

    # 2. Remove Section Numbers (e.g. "1.2", "3.1.5")
    text = re.sub(r'^\s*\d+(\.\d+)+\s+', '', text, flags=re.MULTILINE)

    # 3. Remove Footer Contents
    text = re.sub(r'\n\s*Page\s*\d+\s*of\s*\d+\s*\n', '\n', text, flags=re.IGNORECASE)
    text = re.sub(r'\n\s*\d+\s*\n', '\n', text) 

    # 4. Remove Captions
    text = re.sub(r'^\s*(Figure|Fig\.|Table|Chart|Diagram)\s*\d+.*$', '', text, flags=re.MULTILINE | re.IGNORECASE)

    # 5. Remove Code Blocks
    text = re.sub(r'^.*[{};=<>].*$', '', text, flags=re.MULTILINE) 
    text = re.sub(r'^\s*(def|class|import|from|return|if\s*\(|for\s*\(|while\s*\().*$', '', text, flags=re.MULTILINE)

    # 6. Remove Symbols & Emojis
    text = re.sub(r'[^\w\s.,!?\'"()-]', '', text)

    # 7. Remove Numbered Lists (1., 2., etc.)
    text = re.sub(r'^\s*\d+\.\s*', '', text, flags=re.MULTILINE)

    def roman_to_int(roman):
        values = {'I': 1, 'V': 5, 'X': 10, 'L': 50, 'C': 100}
        total = 0
        prev = 0
        for ch in reversed(roman):
            curr = values.get(ch, 0)
            if curr < prev:
                total -= curr
            else:
                total += curr
                prev = curr
        return total

    # 8. Handle Roman Numerals
    text = re.sub(
        r'(UNIT|CHAPTER)\s*([IVXLCDM]+)\b',
        lambda m: f"{m.group(1).upper()} {roman_to_int(m.group(2).upper())} ",
        text,
        flags=re.IGNORECASE
    )

    # 9. Handle Standard Numerals
    text = re.sub(
        r'(UNIT|CHAPTER)\s*(\d+)\b',
        lambda m: f"{m.group(1).upper()} {m.group(2)} ",
        text,
        flags=re.IGNORECASE
    )

    text = re.sub(r'\[\d+\]', '', text)
    text = re.sub(r'\s+', ' ', text)
    return text.strip()

# --- HELPER FOR BINARY FILES (.ppt) ---
def extract_text_from_binary(filepath):
    # Extracts readable strings from binary files using a robust "strings" like method
    text = ""
    try:
        with open(filepath, 'rb') as f:
            data = f.read()
            
            # 1. Regex to find sequences of 4+ printable ASCII characters (including spaces)
            # This mimics the unix 'strings' command but allows spaces for sentences
            ascii_matches = re.findall(b'[a-zA-Z0-9\s\.,;:!?\-\(\)\[\]]{4,}', data)
            for match in ascii_matches:
                try:
                    # Decode and verify it has at least some letters (avoid pure numbers/symbols)
                    decoded = match.decode('ascii', errors='ignore')
                    if re.search(r'[a-zA-Z]', decoded): 
                        text += decoded + "\n"
                except: pass
            
            # 2. Extract UTF-16LE strings (common in older Office binary formats)
            # Looks for sequences like "T\x00e\x00x\x00t\x00"
            utf16_matches = re.findall(b'(?:[\x20-\x7E]\x00){4,}', data)
            for match in utf16_matches:
                try:
                    decoded = match.decode('utf-16le', errors='ignore')
                    if re.search(r'[a-zA-Z]', decoded):
                        text += decoded + "\n"
                except: pass
                
    except Exception as e:
        print(f"Binary Extraction Error: {e}")
    
    return text

def extract_text(filepath):
    ext = filepath.split('.')[-1].lower()
    text = ""
    try:
        if ext == 'pdf':
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text() or ''
                    lines = page_text.split('\n')
                    if len(lines) > 5: text += '\n'.join(lines[2:-2]) + '\n'
                    else: text += '\n'.join(lines) + '\n'
        elif ext == 'docx':
            doc = Document(filepath)
            for para in doc.paragraphs: text += para.text + '\n'
        elif ext == 'txt':
            with open(filepath, 'r', encoding='utf-8') as f: text = f.read()
        elif ext in ['pptx', 'ppt']:
            # Attempt to use python-pptx (Works for .pptx)
            extracted = False
            if Presentation:
                try:
                    prs = Presentation(filepath)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    if len(text.strip()) > 10:
                        extracted = True
                except Exception:
                    # If python-pptx fails, we fall through to binary extraction
                    pass
            
            # Fallback to binary extraction if library failed or produced no text
            if not extracted:
                print("Falling back to binary extraction for PPT...")
                text = extract_text_from_binary(filepath)

    except Exception:
        return ""
    
    # APPLY CLEANING TO EVERYTHING
    return clean_text(text)

def segment_text_by_units(text):
    pattern = r'(UNIT\s*\d+|CHAPTER\s*\d+)'
    parts = re.split(pattern, text, flags=re.IGNORECASE)
    segments = {}
    current_unit = "General Overview"
    if len(parts) < 2: return {current_unit: text}
    for part in parts:
        part = part.strip()
        if not part: continue
        if re.match(pattern, part, flags=re.IGNORECASE):
            current_unit = part.upper()
            segments[current_unit] = ""
        else:
            if current_unit not in segments:
                segments[current_unit] = ""
            segments[current_unit] += ' ' + part
    return segments

def ml_refine_sentences(sentences, full_text):
    if ml_model is None:
        return [str(s).strip() for s in sentences]

    try:
        texts = [str(s).strip() for s in sentences]
        doc_emb = ml_model.encode(full_text, convert_to_tensor=True)
        sent_emb = ml_model.encode(texts, convert_to_tensor=True)
        scores = util.cos_sim(sent_emb, doc_emb).squeeze()
        ranked = sorted(zip(texts, scores), key=lambda x: x[1], reverse=True)
        refined = [text for text,_ in ranked if len(text.split())>=10 and any(tok.pos_=='VERB' for tok in nlp(text))]
        return refined
    except:
        return [str(s).strip() for s in sentences]

def generate_summary_block(text, style='paragraph', sentence_count=5):
    if len(text.split())<30: return text
    parser = PlaintextParser.from_string(text, Tokenizer("english"))
    summarizer = LexRankSummarizer() if style=='bullets' else LsaSummarizer()
    summarizer.stop_words = get_stop_words("english")
    try:
        sentences = summarizer(parser.document, sentence_count*2)
        final_sentences = ml_refine_sentences(sentences, text)[:sentence_count]
        return final_sentences if style=='bullets' else ' '.join(final_sentences)
    except:
        return text

def translate_summary_dict(summary_dict, target_lang):
    if target_lang=='en': return summary_dict
    translated = {}
    translator = GoogleTranslator(source='auto', target=target_lang)
    for k,v in summary_dict.items():
        try:
            t_title = translator.translate(k)
            if isinstance(v, list): t_content = [translator.translate(p) for p in v]
            else: t_content = translator.translate(str(v))
            translated[t_title] = t_content
        except: translated[k] = v
        time.sleep(0.5)
    return translated

def sanitize_for_pdf(text):
    return text.encode('latin-1', 'replace').decode('latin-1')

def analyze_document_light(text):
    doc = nlp(text[:100000])
    
    candidates = []
    for chunk in doc.noun_chunks:
        clean = chunk.text.lower().strip()
        if len(clean) > 3 and not chunk.root.is_stop and not chunk.root.is_punct:
            candidates.append(clean)
    top_keywords = [item[0].title() for item in Counter(candidates).most_common(10)]

    entities = defaultdict(set)
    for ent in doc.ents:
        val = ent.text.strip()
        if len(val) < 2: continue
        if ent.label_ == 'ORG': entities['Organizations'].add(val)
        elif ent.label_ == 'DATE': entities['Dates'].add(val)
        
    final_entities = {k: sorted(list(v))[:8] for k, v in entities.items() if v}

    return {
        'word_count': len(text.split()), 
        'keywords': top_keywords,
        'entities': final_entities 
    }
# --- ROUTES ---

@app.route('/')
def home():
    if is_maintenance_active():
        note = Notification.query.order_by(Notification.created_at.desc()).first()
        custom_msg = note.message if note else None
        return render_template('maintenance.html', message=custom_msg)
    return render_template('home.html')

@app.route('/login', methods=['GET','POST'])
def login():
    if request.method=='POST':
        user = User.query.filter((User.username==request.form.get('login_id'))|(User.email==request.form.get('login_id'))).first()
        if user:
            if user.is_blocked:
                flash('Your account has been blocked.')
                return render_template('login.html')
            if check_password_hash(user.password, request.form.get('password')):
                login_user(user)
                log_activity('login_success')
                if user.role == 'admin':
                    return redirect(url_for('admin_dashboard'))
                return redirect(url_for('dashboard'))
        flash('Invalid Credentials')
    return render_template('login.html')

@app.route('/signup', methods=['GET','POST'])
def signup():
    if request.method=='POST':
        hashed_pw = generate_password_hash(request.form.get('password'), method='pbkdf2:sha256')
        role = 'admin' if User.query.count() == 0 else 'user'
        new_user = User(username=request.form.get('username'), email=request.form.get('email'), dob=request.form.get('dob'), password=hashed_pw, role=role)
        db.session.add(new_user); db.session.commit()
        return redirect(url_for('login'))
    return render_template('signup.html')

@app.route('/logout')
@login_required
def logout():
    log_activity('logout'); logout_user()
    return redirect(url_for('login'))

@app.route('/dashboard', methods=['GET', 'POST'])
@login_required
def dashboard():
    if is_maintenance_active():
        return redirect(url_for('home'))
    if request.method == 'POST':
        file = request.files.get('file'); raw_text = request.form.get('raw_text')
        text, filename = "", ""
        if file and file.filename:
            filename = secure_filename(file.filename); save_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(save_path); text = extract_text(save_path); log_activity('upload_doc', filename)
        elif raw_text:
            text = clean_text(raw_text); filename = f"Text_{datetime.now().strftime('%Y%m%d')}.txt"; log_activity('upload_text')
            
        if not text: flash("No text found."); return redirect(url_for('dashboard'))
        analysis = analyze_document_light(text); unit_segments = segment_text_by_units(text)
        style = request.form.get('style'); length_opt = request.form.get('length'); language = request.form.get('language')
        target = 20 if length_opt=='short' else 120 if length_opt=='long' else 60
        sent_per_unit = max(4, int(target / max(1, len(unit_segments))))
        
        # --- UPDATED: Sort units numerically ---
        def get_unit_num(key):
            match = re.search(r'\d+', key)
            return int(match.group()) if match else 0

        sorted_keys = sorted(unit_segments.keys(), key=get_unit_num)
        summary_dict = {}
        for t in sorted_keys:
            content = unit_segments[t]
            summary = generate_summary_block(content, style, sent_per_unit)
            if summary:
                summary_dict[t] = summary
        
        if language and language != 'en': summary_dict = translate_summary_dict(summary_dict, language)
        new_sum = Summary(filename=filename, original_text=text, summary_text=json.dumps(summary_dict), language=language, user_id=current_user.id)
        db.session.add(new_sum); db.session.commit()
        return render_template('result.html', summary_dict=summary_dict, analysis=analysis, summary_id=new_sum.id, filename=filename, language=language, original=text, chat_history=[])
    return render_template('dashboard.html')

# --- PDF & IMAGE TOOLS ---

@app.route('/tools/split_pdf', methods=['GET', 'POST'])
@login_required
def split_pdf_tool():
    if request.method == 'POST':
        file = request.files.get('file')
        if not file: return redirect(request.url)
        try:
            log_activity('split_pdf')
            reader = PyPDF2.PdfReader(file); writer = PyPDF2.PdfWriter()
            page_range = request.form.get('page_range', '')
            total = len(reader.pages); selected = set()
            for part in page_range.split(','):
                if '-' in part:
                    s, e = map(int, part.split('-'))
                    for p in range(s, e+1): 
                        if 1 <= p <= total: selected.add(p-1)
                elif part.strip().isdigit():
                    p = int(part.strip())
                    if 1 <= p <= total: selected.add(p-1)
            for p_num in sorted(selected): writer.add_page(reader.pages[p_num])
            out = io.BytesIO(); writer.write(out); out.seek(0)
            return send_file(out, as_attachment=True, download_name=f"Split_{file.filename}", mimetype='application/pdf')
        except: flash("Error processing PDF"); return redirect(request.url)
    return render_template('split_pdf.html')

@app.route('/tools/merge_pdf', methods=['GET', 'POST'])
@login_required
def merge_pdf_tool():
    if request.method == 'POST':
        files = request.files.getlist('files')
        if not files or files[0].filename == '': return redirect(request.url)
        try:
            log_activity('merge_pdf')
            merger = PyPDF2.PdfMerger()
            for f in files: merger.append(f)
            out = io.BytesIO(); merger.write(out); merger.close(); out.seek(0)
            return send_file(out, as_attachment=True, download_name="Merged_Document.pdf", mimetype='application/pdf')
        except: flash("Merge Error"); return redirect(request.url)
    return render_template('merge_pdf.html')

@app.route('/tools/img_to_pdf', methods=['GET', 'POST'])
@login_required
def img_to_pdf_tool():
    if request.method == 'POST':
        # Check for files
        if 'files' not in request.files:
            flash('No files uploaded')
            return redirect(request.url)
        
        files = request.files.getlist('files')
        valid_files = [f for f in files if f.filename != '']
        
        if not valid_files:
            flash('No images selected')
            return redirect(request.url)

        # Get Options
        page_size = request.form.get('page_size', 'original')
        orientation = request.form.get('orientation', 'portrait')
        margin_opt = request.form.get('margin', 'none')
        is_preview = request.form.get('is_preview') == 'true'  # Check for Preview Flag

        # Define Sizes & Margins
        SIZES = {'a4': (1240, 1754), 'letter': (1275, 1650)}
        MARGINS = {'none': 0, 'small': 40, 'big': 100}
        margin_px = MARGINS.get(margin_opt, 0)

        try:
            processed_images = []
            
            for file in valid_files:
                img = Image.open(file)
                if img.mode in ('RGBA', 'P'):
                    bg = Image.new("RGB", img.size, (255, 255, 255))
                    if img.mode == 'RGBA': bg.paste(img, mask=img.split()[3])
                    else: img = img.convert('RGB')
                    img = bg

                if page_size == 'original':
                    processed_images.append(img)
                else:
                    page_w, page_h = SIZES.get(page_size, SIZES['a4'])
                    if orientation == 'landscape': page_w, page_h = page_h, page_w
                    
                    canvas = Image.new('RGB', (page_w, page_h), (255, 255, 255))
                    safe_w, safe_h = page_w - (2 * margin_px), page_h - (2 * margin_px)
                    
                    if safe_w <= 0: safe_w, safe_h = page_w, page_h

                    img_ratio = img.width / img.height
                    target_ratio = safe_w / safe_h
                    
                    if img_ratio > target_ratio:
                        new_w = safe_w
                        new_h = int(safe_w / img_ratio)
                    else:
                        new_h = safe_h
                        new_w = int(safe_h * img_ratio)
                        
                    img_resized = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
                    x_off = margin_px + (safe_w - new_w) // 2
                    y_off = margin_px + (safe_h - new_h) // 2
                    
                    canvas.paste(img_resized, (x_off, y_off))
                    processed_images.append(canvas)
            
            if processed_images:
                output_buffer = io.BytesIO()
                processed_images[0].save(output_buffer, save_all=True, append_images=processed_images[1:], format='PDF')
                output_buffer.seek(0)

                # Return Inline for Preview, Attachment for Download
                return send_file(
                    output_buffer,
                    as_attachment=not is_preview, 
                    download_name=f"Images_Converted_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                    mimetype='application/pdf'
                )

        except Exception as e:
            print(f"Error: {e}")
            flash(f"Error converting images: {str(e)}")
            return redirect(request.url)

    return render_template('img_to_pdf.html')

# --- ADMIN PANEL ROUTES ---

@app.route('/admin')
@admin_required
def admin_dashboard():
    u_count = User.query.count(); d_count = Summary.query.count()
    s_count = ActivityLog.query.filter_by(action='split_pdf').count()
    m_count = ActivityLog.query.filter_by(action='merge_pdf').count()
    i_count = ActivityLog.query.filter_by(action='img_to_pdf').count()
    logs = ActivityLog.query.order_by(ActivityLog.timestamp.desc()).limit(10).all()
    ratings = Rating.query.join(User).add_columns(User.username, Rating.stars, Rating.timestamp).order_by(Rating.timestamp.desc()).limit(20).all()
    avg = db.session.query(db.func.avg(Rating.stars)).scalar() or 0
    return render_template('admin_dashboard.html', user_count=u_count, doc_count=d_count, split_count=s_count, merge_count=m_count, img_count=i_count, logs=logs, ratings=ratings, avg_rating=round(avg, 1))

@app.route('/admin/users', methods=['GET', 'POST'])
@admin_required
def admin_users():
    if request.method == 'POST':
        action = request.form.get('action'); uid = request.form.get('user_id'); user = User.query.get(uid)
        if user:
            if action == 'edit_user':
                user.username = request.form.get('username'); user.email = request.form.get('email'); flash("Updated.")
            elif user.id != current_user.id:
                if action == 'block': user.is_blocked = True
                elif action == 'unblock': user.is_blocked = False
                elif action == 'make_admin': user.role = 'admin'
                elif action == 'delete':
                    Summary.query.filter_by(user_id=user.id).delete()
                    db.session.delete(user)
            db.session.commit()
    return render_template('admin_users.html', users=User.query.all())

@app.route('/admin/settings', methods=['POST'])
@admin_required
def admin_settings():
    # Handle Maintenance Mode
    m_mode = 'true' if request.form.get('maintenance_mode') else 'false'
    s_maint = SiteSetting.query.filter_by(key_name='maintenance_mode').first()
    if not s_maint: s_maint = SiteSetting(key_name='maintenance_mode'); db.session.add(s_maint)
    s_maint.value = m_mode
    
    # Handle Ratings Toggle
    r_mode = 'true' if request.form.get('enable_ratings') else 'false'
    s_rate = SiteSetting.query.filter_by(key_name='enable_ratings').first()
    if not s_rate: 
        s_rate = SiteSetting(key_name='enable_ratings', value='false')
        db.session.add(s_rate)
    
    # NEW LOGIC: If enabling (transition from False -> True), reset the session timer
    if r_mode == 'true' and s_rate.value != 'true':
        s_time = SiteSetting.query.filter_by(key_name='rating_start_time').first()
        if not s_time:
            s_time = SiteSetting(key_name='rating_start_time')
            db.session.add(s_time)
        s_time.value = datetime.utcnow().isoformat()

    s_rate.value = r_mode

    # Handle Announcement
    ann = request.form.get('announcement')
    if ann: db.session.add(Notification(message=ann))
    
    db.session.commit(); flash("Settings updated."); return redirect(url_for('admin_dashboard'))

# --- CORE USER ROUTES ---

@app.route('/profile', methods=['GET', 'POST'])
@login_required
def profile():
    if request.method == 'POST':
        u = request.form.get('username'); e = request.form.get('email')
        exist = User.query.filter(((User.username == u) | (User.email == e)) & (User.id != current_user.id)).first()
        if exist:
            flash('Username or Email already taken.', 'error')
        else:
            current_user.username = u; current_user.email = e
            db.session.commit(); flash('Profile updated successfully!', 'success')
    return render_template('user_profile.html')

@app.route('/history')
@login_required
def history():
    if is_maintenance_active(): return redirect(url_for('home'))
    sums = Summary.query.filter_by(user_id=current_user.id).order_by(Summary.created_at.desc()).all()
    return render_template('history.html', summaries=sums)

@app.route('/view/<int:id>')
@login_required
def view_summary(id):
    s = Summary.query.get_or_404(id)
    if current_user.role != 'admin' and s.user_id != current_user.id: return redirect(url_for('history'))
    try: s_dict = json.loads(s.summary_text)
    except: s_dict = {"Summary": s.summary_text}
    
    # --- FIX: Retrieve Chat History from DB ---
    # Previously this was passed as empty list []
    history = ChatMessage.query.filter_by(summary_id=id).order_by(ChatMessage.timestamp.asc()).all()
    
    return render_template('result.html', summary_dict=s_dict, analysis=analyze_document_light(s.original_text), summary_id=s.id, filename=s.filename, language=s.language, original=s.original_text, chat_history=history)

@app.route('/chat_api', methods=['POST'])
@login_required
def chat_api():
    data = request.json
    text = data.get('context', '')
    question = data.get('question', '').strip()
    summary_id = data.get('summary_id') 
    
    # Clean text
    text = re.sub(r'\s+', ' ', text)
    sentences = re.split(r'(?<!\w\.\w.)(?<![A-Z][a-z]\.)(?<=\.|\?)\s', text)
    
    answer = None
    best_idx = -1

    # METHOD 1: Try Deep Learning (Semantic Search) first if model is loaded
    if ml_model:
        try:
            # Encode the question and all sentences into vectors
            sentence_embeddings = ml_model.encode(sentences, convert_to_tensor=True)
            question_embedding = ml_model.encode(question, convert_to_tensor=True)

            # Compute similarity scores
            # util.cos_sim returns a matrix, we want the first row
            scores = util.cos_sim(question_embedding, sentence_embeddings)[0]

            # Find the sentence with the highest score
            best_idx = int(scores.argmax())
            best_score = float(scores[best_idx])

            # Threshold: Only accept if similarity is > 0.25 (adjustable)
            if best_score < 0.25:
                best_idx = -1
        except Exception as e:
            print(f"ML Search Error: {e}")
            best_idx = -1

    # METHOD 2: Fallback to TF-IDF if ML fails or finds nothing specific
    if best_idx == -1:
        try:
            corpus = sentences + [question]
            vectorizer = TfidfVectorizer(stop_words='english')
            tfidf_matrix = vectorizer.fit_transform(corpus)
            cosine_sims = cosine_similarity(tfidf_matrix[-1], tfidf_matrix[:-1])
            best_idx = cosine_sims.argmax()
            if cosine_sims[0, best_idx] < 0.1: # Low threshold for keyword match
                best_idx = -1
        except:
            best_idx = -1

    if best_idx != -1:
        # Return context: Previous sentence + Best Match + Next Sentence
        start = max(0, best_idx)
        end = min(len(sentences), best_idx + 3)
        
        # Formatting: Clean the selected sentences
        selected_sentences = sentences[start:end]
        cleaned_sentences = []
        for s in selected_sentences:
            s = re.sub(r'^[\d\W_]+', '', s).strip() # aggressive strip of bullets
            if s:
                cleaned_sentences.append(s[0].upper() + s[1:])

        answer = " ".join(cleaned_sentences)
    else:
        answer = "I couldn't find a specific answer in the document."

    if summary_id:
        db.session.add(ChatMessage(summary_id=summary_id, user_id=current_user.id, message=question, sender='user'))
        db.session.add(ChatMessage(summary_id=summary_id, user_id=current_user.id, message=answer, sender='bot'))
        db.session.commit()
        
    return jsonify({'answer': answer})

@app.route('/forgot_password', methods=['GET', 'POST'])
def forgot_password():
    if request.method == 'POST':
        email = request.form.get('email'); dob = request.form.get('dob'); new_pw = request.form.get('new_password')
        if len(new_pw) < 8: flash('Password must be 8+ chars.'); return redirect(url_for('forgot_password'))
        user = User.query.filter_by(email=email, dob=dob).first()
        if user:
            user.password = generate_password_hash(new_pw, method='pbkdf2:sha256'); db.session.commit()
            flash('Success! Please login.'); return redirect(url_for('login'))
        flash('Verification failed.')
    return render_template('forgot_password.html')

@app.route('/delete/<int:id>')
@login_required
def delete_summary(id):
    summary = Summary.query.get_or_404(id)
    if summary.user_id == current_user.id or current_user.role == 'admin':
        ChatMessage.query.filter_by(summary_id=id).delete()
        db.session.delete(summary); db.session.commit()
    return redirect(url_for('admin_dashboard' if current_user.role == 'admin' else 'history'))

# --- DOWNLOAD ROUTES (RESTORED) ---

@app.route('/download_pdf/<int:id>')
@login_required
def download_pdf(id):
    summary = Summary.query.get_or_404(id)
    if summary.language and summary.language != 'en':
        flash("PDF Download is restricted to English summaries.")
        return redirect(url_for('view_summary', id=id))
    
    summary_data = json.loads(summary.summary_text)
    analysis = analyze_document_light(summary.original_text)
    
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    pdf.cell(0, 10, txt=f"Summary: {sanitize_for_pdf(summary.filename)}", ln=True, align='C')
    pdf.ln(5)
    
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(0, 10, txt="Key Insights", ln=True)
    pdf.set_font("Arial", '', 10)
    if analysis['keywords']:
        pdf.multi_cell(0, 6, txt="Keywords: " + ", ".join(analysis['keywords']))
    pdf.ln(5)
    
    for title, content in summary_data.items():
        pdf.set_font("Arial", 'B', 14)
        pdf.cell(0, 10, txt=sanitize_for_pdf(title), ln=True)
        pdf.set_font("Arial", size=11)
        if isinstance(content, list):
            for para in content:
                pdf.multi_cell(0, 6, txt=sanitize_for_pdf(para))
                pdf.ln(4)
        else:
            pdf.multi_cell(0, 6, txt=sanitize_for_pdf(str(content)))
        pdf.ln(5)
        
    response = make_response(pdf.output(dest='S').encode('latin-1'))
    response.headers['Content-Type'] = 'application/pdf'
    response.headers['Content-Disposition'] = f'attachment; filename=summary_{id}.pdf'
    return response

@app.route('/download_word/<int:id>')
@login_required
def download_word(id):
    summary = Summary.query.get_or_404(id)
    summary_data = json.loads(summary.summary_text)
    analysis = analyze_document_light(summary.original_text)
    
    doc = Document()
    doc.add_heading(f"Summary: {summary.filename}", 0)
    doc.add_heading("Key Insights", level=1)
    if analysis['keywords']:
        doc.add_paragraph("Keywords: " + ", ".join(analysis['keywords']))
    
    for title, content in summary_data.items():
        doc.add_heading(title, level=1)
        if isinstance(content, list):
            for para in content:
                doc.add_paragraph(para, style='List Bullet')
        else:
            doc.add_paragraph(str(content))
            
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    return send_file(file_stream, as_attachment=True, download_name=f"summary_{id}.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

# --- NEW BULK DELETE ROUTE (ADDED) ---
@app.route('/delete_bulk', methods=['POST'])
@login_required
def delete_bulk():
    data = request.json
    ids = data.get('ids', [])
    if ids:
        # Delete associated chats first
        ChatMessage.query.filter(ChatMessage.summary_id.in_(ids), ChatMessage.user_id == current_user.id).delete(synchronize_session=False)
        # Delete summaries
        Summary.query.filter(Summary.id.in_(ids), Summary.user_id == current_user.id).delete(synchronize_session=False)
        db.session.commit()
    return jsonify({'status': 'success'})

@app.route('/submit_rating_api', methods=['POST'])
@login_required
def submit_rating_api():
    data = request.json
    stars = int(data.get('rating', 0))
    if 1 <= stars <= 5:
        new_rating = Rating(user_id=current_user.id, stars=stars)
        db.session.add(new_rating)
        db.session.commit()
        return jsonify({'status': 'success'})
    return jsonify({'status': 'error'}), 400

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
        # Initialize default settings if missing
        if not SiteSetting.query.filter_by(key_name='maintenance_mode').first():
            db.session.add(SiteSetting(key_name='maintenance_mode', value='false'))
        if not SiteSetting.query.filter_by(key_name='enable_ratings').first():
            db.session.add(SiteSetting(key_name='enable_ratings', value='false'))
        db.session.commit()
    app.run(debug=True)